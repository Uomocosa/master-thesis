from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu

from create_slides.UnisiBackground import UnisiBackground


def apply_unisi_background(pptx_slide: PptxSlide, background: UnisiBackground) -> None:
    if background.banner_image is not None:
        banner = pptx_slide.shapes.add_picture(
            str(background.banner_image),
            Emu(background.banner_position[0]),
            Emu(background.banner_position[1]),
            Emu(background.banner_size[0]),
            Emu(background.banner_size[1]),
        )
        # z-order is spTree element order; the layout placeholders (title/body)
        # already exist, so move the banner behind them (index 2 = first shape
        # slot after nvGrpSpPr/grpSpPr).
        sp_tree = pptx_slide.shapes._spTree
        sp_tree.remove(banner._element)
        sp_tree.insert(2, banner._element)
    pptx_slide.shapes.add_picture(
        str(background.logo_image),
        Emu(background.logo_position[0]),
        Emu(background.logo_position[1]),
        Emu(background.logo_size[0]),
        Emu(background.logo_size[1]),
    )

    line = background.line_separator
    x, y = line.position
    connector = pptx_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Emu(x), Emu(y), Emu(x), Emu(y + line.length)
    )
    connector.line.color.rgb = RGBColor.from_string(line.color)
    connector.line.width = Emu(line.weight)
