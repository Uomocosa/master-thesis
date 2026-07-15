from pptx import Presentation
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu

from create_slides.DeckBuilding.AddSlideImagesMethod import add_slide_images
from create_slides.DeckBuilding.AddSlideMoviesMethod import add_slide_movies
from create_slides.DeckBuilding.ApplyUnisiBackgroundMethod import apply_unisi_background
from create_slides.DeckBuilding.ApplyUnisiTextStyleMethod import apply_unisi_text_style
from create_slides.Slide import Slide
from create_slides.UnisiBackground import UnisiBackground
from create_slides.UnisiText import UnisiText
from create_slides.UnisiTitle import UnisiTitle


def add_unisi_slide(
    prs: Presentation, slide: Slide, background: UnisiBackground | None = None
) -> PptxSlide:
    layout = prs.slide_layouts[1]
    pptx_slide = prs.slides.add_slide(layout)
    if background is not None:
        apply_unisi_background(pptx_slide, background)

    pptx_slide.shapes.title.text = slide.title
    if background is not None:
        title_shape = pptx_slide.shapes.title
        title_shape.left, title_shape.top = (Emu(v) for v in background.title_position)
        title_shape.width, title_shape.height = (Emu(v) for v in background.title_size)
        apply_unisi_text_style(title_shape.text_frame.paragraphs[0], UnisiTitle())

    body_shape = pptx_slide.placeholders[1]
    body = body_shape.text_frame
    body.clear()
    body.text = slide.visual_bullets[0] if slide.visual_bullets else ""
    for bullet in slide.visual_bullets[1:]:
        paragraph = body.add_paragraph()
        paragraph.text = bullet
    if background is not None:
        body_x, body_y = background.body_position
        body_w, body_h = background.body_size
        body_shape.left, body_shape.top = Emu(body_x), Emu(body_y)
        body_shape.height = Emu(body_h)
        if slide.movies:
            # Bullets on top, click-to-play videos fill the rest of the body.
            text_h = int(body_h * 0.30)
            body_shape.width = Emu(body_w)
            body_shape.height = Emu(text_h)
            movies_y = body_y + int(body_h * 0.34)
            add_slide_movies(
                pptx_slide,
                slide.movies,
                (body_x, movies_y),
                (body_w, body_h - int(body_h * 0.34)),
            )
        elif slide.images and slide.images_below:
            # Stacked layout: full-width bullets on top, images fill the rest.
            text_h = int(body_h * 0.42)
            body_shape.width = Emu(body_w)
            body_shape.height = Emu(text_h)
            images_y = body_y + int(body_h * 0.46)
            add_slide_images(
                pptx_slide,
                slide.images,
                (body_x, images_y),
                (body_w, body_h - int(body_h * 0.46)),
            )
        elif slide.images:
            # Two-column layout: bullets on the left, images fill the right.
            text_w = int(body_w * 0.52)
            body_shape.width = Emu(text_w)
            image_x = body_x + int(body_w * 0.56)
            add_slide_images(
                pptx_slide,
                slide.images,
                (image_x, body_y),
                (body_w - int(body_w * 0.56), body_h),
            )
        else:
            body_shape.width = Emu(body_w)
        for paragraph in body.paragraphs:
            apply_unisi_text_style(paragraph, UnisiText.level(1))

    notes = slide.speech
    if slide.todos:
        notes += "\n\nTODO:\n" + "\n".join(f"- {todo}" for todo in slide.todos)
    pptx_slide.notes_slide.notes_text_frame.text = notes

    return pptx_slide
