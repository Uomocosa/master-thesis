from pathlib import Path

from pptx.parts.image import Image
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu

GAP = Emu(228600)  # 0.25 in between images


def add_slide_images(
    pptx_slide: PptxSlide,
    images: list[Path],
    area_position: tuple[int, int],
    area_size: tuple[int, int],
) -> None:
    """Lay images out inside the area, centered as a group.

    Wide areas get one row (equal heights, widths proportional to aspect);
    tall areas get one column (equal widths, heights proportional to aspect).
    """
    if not images:
        return

    area_x, area_y = area_position
    area_w, area_h = area_size

    aspects = []
    for image_path in images:
        px_w, px_h = Image.from_file(str(image_path)).size
        aspects.append(px_w / px_h)

    if area_w >= area_h:
        usable_w = area_w - GAP * (len(images) - 1)
        row_h = min(area_h, int(usable_w / sum(aspects)))
        row_w = int(row_h * sum(aspects)) + GAP * (len(images) - 1)
        left = area_x + (area_w - row_w) // 2
        top = area_y + (area_h - row_h) // 2
        for image_path, aspect in zip(images, aspects):
            img_w = int(row_h * aspect)
            pptx_slide.shapes.add_picture(
                str(image_path), Emu(left), Emu(top), Emu(img_w), Emu(row_h)
            )
            left += img_w + GAP
    else:
        usable_h = area_h - GAP * (len(images) - 1)
        col_w = min(area_w, int(usable_h / sum(1 / a for a in aspects)))
        col_h = int(col_w * sum(1 / a for a in aspects)) + GAP * (len(images) - 1)
        left = area_x + (area_w - col_w) // 2
        top = area_y + (area_h - col_h) // 2
        for image_path, aspect in zip(images, aspects):
            img_h = int(col_w / aspect)
            pptx_slide.shapes.add_picture(
                str(image_path), Emu(left), Emu(top), Emu(col_w), Emu(img_h)
            )
            top += img_h + GAP
