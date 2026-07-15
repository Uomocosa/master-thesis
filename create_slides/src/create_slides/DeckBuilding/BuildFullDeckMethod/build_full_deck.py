from pathlib import Path

from loguru import logger
from pptx import Presentation

from create_slides.DeckBuilding.AddUnisiSlideMethod import add_unisi_slide
from create_slides.DeckBuilding.ClearTemplateSlidesMethod import clear_template_slides
from create_slides.Slide import Slide
from create_slides.UnisiBackground import UnisiBackground


def build_full_deck(
    slides: list[tuple[Slide, UnisiBackground | None]], output_path: Path
) -> Path:
    template = next((bg for _, bg in slides if bg is not None), None)
    if template is not None:
        prs = Presentation(str(template.template_path))
        clear_template_slides(prs)
    else:
        prs = Presentation()

    for slide, background in slides:
        add_unisi_slide(prs, slide, background)

    # Mark hidden slides: skipped in slideshow, still present/visible in edit view.
    # `show` lives on the <p:sld> root of the slide part; PowerPoint ignores it on <p:sldId>.
    for (slide, _), built in zip(slides, prs.slides):
        if slide.hidden:
            built._element.set("show", "0")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"Saved {len(slides)} slides to {output_path}")
    return output_path
