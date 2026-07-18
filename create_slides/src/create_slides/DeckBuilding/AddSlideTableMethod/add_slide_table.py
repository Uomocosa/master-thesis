from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu, Pt

from create_slides.SlideTable import SlideTable

UNISI_RED = RGBColor(0x99, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
FONT_NAME = "Optima"
HEADER_SIZE_PT = 18
BODY_SIZE_PT = 16


def _style_cell(cell, text: str, *, header: bool, align: PP_ALIGN) -> None:
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = UNISI_RED if header else WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = FONT_NAME
        run.font.size = Pt(HEADER_SIZE_PT if header else BODY_SIZE_PT)
        run.font.bold = header
        run.font.color.rgb = WHITE if header else BLACK


def add_slide_table(
    pptx_slide: PptxSlide,
    table: SlideTable,
    area_position: tuple[int, int],
    area_size: tuple[int, int],
) -> None:
    x, y = area_position
    w, h = area_size
    n_rows = 1 + len(table.rows)
    n_cols = len(table.headers)
    frame = pptx_slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(w), Emu(h))
    pptx_table = frame.table
    for col, header in enumerate(table.headers):
        _style_cell(pptx_table.cell(0, col), header, header=True, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(table.rows, start=1):
        for col, value in enumerate(row):
            # First column reads as a label, the rest as values.
            align = PP_ALIGN.LEFT if col == 0 else PP_ALIGN.CENTER
            _style_cell(pptx_table.cell(row_index, col), value, header=False, align=align)
