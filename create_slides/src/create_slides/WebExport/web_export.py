"""Export an edited .pptx to a static, self-contained reveal.js web presentation.

The pptx is the source of truth (it may have been hand-edited after generation).
Each slide is rendered to a full PNG via PowerPoint COM (``render_pptx.ps1``) and used
as a reveal.js slide background; embedded videos are extracted and overlaid as
absolutely-positioned ``<video>`` elements at the exact rect of their pptx shape. Videos
auto-play when their slide is shown; a second video on the same slide starts 5 s later.
"""

import html
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Relationship-namespace attribute keys used by python-pptx lxml elements.
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_EMBED = f"{{{_R_NS}}}embed"
_LINK = f"{{{_R_NS}}}link"


@dataclass
class VideoOverlay:
    """One embedded video and its position as fractions of the slide (0..1)."""

    src: str
    left: float
    top: float
    width: float
    height: float


@dataclass
class WebSlide:
    index: int
    image: str
    videos: list[VideoOverlay] = field(default_factory=list)


def _find_video_blob(shape, slide_part) -> bytes | None:
    """Resolve the embedded mp4 bytes for a media shape via its r:embed/r:link rels."""
    rels = slide_part.rels
    for attr in (_EMBED, _LINK):
        for el in shape._element.iter():
            rid = el.get(attr)
            if not rid or rid not in rels:
                continue
            part = rels[rid].target_part
            ctype = getattr(part, "content_type", "") or ""
            if "video" in ctype or str(getattr(part, "partname", "")).endswith(".mp4"):
                return part.blob
    return None


def hidden_slide_numbers(pptx_path: Path) -> set[int]:
    """1-based positions of slides hidden from the slideshow (<p:sld show="0">)."""
    prs = Presentation(str(pptx_path))
    return {
        i
        for i, slide in enumerate(prs.slides, start=1)
        if slide._element.get("show") in ("0", "false")
    }


def extract_slides(
    pptx_path: Path, out_assets: Path, image_names: dict[int, str], hidden: set[int]
) -> list[WebSlide]:
    """Walk the pptx; per slide collect its PNG name + extracted video overlays."""
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    web_slides: list[WebSlide] = []

    for i, slide in enumerate(prs.slides, start=1):
        # Hidden slides are skipped by the PowerPoint slideshow, so skip them here too.
        if i in hidden:
            logger.info(f"slide {i}: hidden in pptx, skipped")
            continue

        videos: list[VideoOverlay] = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.MEDIA:
                continue
            blob = _find_video_blob(shape, slide.part)
            if blob is None:
                logger.warning(f"slide {i}: media shape without resolvable video blob, skipped")
                continue
            name = f"slide{i:02d}_video{len(videos) + 1}.mp4"
            (out_assets / name).write_bytes(blob)
            videos.append(
                VideoOverlay(
                    src=f"assets/{name}",
                    left=shape.left / slide_w,
                    top=shape.top / slide_h,
                    width=shape.width / slide_w,
                    height=shape.height / slide_h,
                )
            )
        web_slides.append(WebSlide(index=i, image=f"assets/{image_names[i]}", videos=videos))
        logger.info(f"slide {i}: {len(videos)} video(s)")
    return web_slides


def render_pngs(
    pptx_path: Path, render_ps1: Path, out_assets: Path, hidden: set[int]
) -> dict[int, str]:
    """Render slides to PNG via PowerPoint COM; copy into assets as slideNN.png.

    COM exports hidden slides too, so they are dropped here rather than shipped as
    unreferenced-but-publicly-fetchable files.
    """
    tmp = out_assets / "_render"
    tmp.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "powershell",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(render_ps1),
            "-PptxPath",
            str(pptx_path),
            "-OutDir",
            str(tmp),
        ],
        check=True,
    )
    names: dict[int, str] = {}
    for png in sorted(tmp.glob("*_slide*.png")):
        num = int(png.stem.split("slide")[-1])
        if num in hidden:
            continue
        dest = f"slide{num:02d}.png"
        shutil.copyfile(png, out_assets / dest)
        names[num] = dest
    shutil.rmtree(tmp, ignore_errors=True)
    logger.info(f"rendered {len(names)} slide PNGs")
    return names


def _slide_section(slide: WebSlide) -> str:
    overlays = "\n".join(
        f'      <video class="vid" data-delay="{j * 5000}" muted playsinline '
        f'preload="auto" '
        f'style="left:{v.left * 100:.4f}%;top:{v.top * 100:.4f}%;'
        f'width:{v.width * 100:.4f}%;height:{v.height * 100:.4f}%">'
        f'<source src="{html.escape(v.src)}" type="video/mp4"></video>'
        for j, v in enumerate(slide.videos)
    )
    return (
        f'    <section data-background-size="contain" data-background-color="#ffffff" '
        f'data-background-image="{html.escape(slide.image)}">\n'
        f'{overlays}\n    </section>'
    )


def render_index_html(slides: list[WebSlide], title: str) -> str:
    sections = "\n".join(_slide_section(s) for s in slides)
    return f"""<!doctype html>
<html lang="it">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{html.escape(title)}</title>
<link rel="stylesheet" href="assets/reveal/reveal.css">
<link rel="stylesheet" href="assets/reveal/theme/white.css">
<style>
  html, body {{ background:#111; }}
  .reveal .slides section {{ padding:0; height:100%; }}
  .reveal .vid {{
    position:absolute; object-fit:contain; background:transparent;
    opacity:0; transition:opacity .25s ease; pointer-events:none;
  }}
  .reveal .vid.visible {{ opacity:1; }}
  .hint {{
    position:fixed; bottom:10px; left:0; right:0; text-align:center;
    font-family:Calibri, Arial, sans-serif; font-size:14px; color:#888;
    z-index:40; pointer-events:none;
  }}
  .fullscreen-btn {{
    position:fixed; top:12px; right:12px; z-index:50;
    width:40px; height:40px; border:none; border-radius:6px;
    background:rgba(0,0,0,.35); color:#fff; font-size:22px; line-height:40px;
    text-align:center; cursor:pointer; opacity:.55; transition:opacity .2s ease;
    padding:0;
  }}
  .fullscreen-btn:hover {{ opacity:1; }}
</style>
</head>
<body>
<div class="reveal">
  <div class="slides">
{sections}
  </div>
</div>
<div class="hint">Clicca per avanzare &middot; le animazioni partono automaticamente &middot; F o &#x26F6; per lo schermo intero</div>
<button class="fullscreen-btn" title="Schermo intero (F)" aria-label="Schermo intero">&#x26F6;</button>
<script src="assets/reveal/reveal.js"></script>
<script>
  Reveal.initialize({{
    controls: true, progress: true, hash: true,
    center: false, width: 1920, height: 1080, margin: 0,
    keyboard: {{ 8: 'prev' }}, transition: 'fade',
  }});

  // Fullscreen toggle: the corner button and reveal's own F shortcut both work.
  function toggleFullscreen() {{
    if (document.fullscreenElement) document.exitFullscreen();
    else document.documentElement.requestFullscreen();
  }}
  document.querySelector('.fullscreen-btn').addEventListener('click', ev => {{
    ev.stopPropagation();
    toggleFullscreen();
  }});

  // Videos auto-play when their slide is shown; each one waits its data-delay ms
  // (0 for the first, 5000 for the second). Leaving the slide cancels pending
  // timers so a not-yet-started video never plays over a different slide.
  function stopVideos(slide) {{
    if (!slide) return;
    slide.querySelectorAll('.vid').forEach(v => {{
      if (v._timer) {{ clearTimeout(v._timer); v._timer = null; }}
      v.classList.remove('visible');
      try {{ v.pause(); }} catch (e) {{}}
    }});
  }}
  function startVideos(slide) {{
    if (!slide) return;
    slide.querySelectorAll('.vid').forEach(v => {{
      v._timer = setTimeout(() => {{
        v._timer = null;
        v.classList.add('visible');
        try {{ v.currentTime = 0; v.play(); }} catch (e) {{}}
      }}, +v.dataset.delay || 0);
    }});
  }}
  Reveal.on('slidechanged', e => {{ stopVideos(e.previousSlide); startVideos(e.currentSlide); }});
  Reveal.on('ready', e => startVideos(e.currentSlide));

  // Click anywhere advances the deck (reveal.js does not do this by default).
  // Guard clicks on the built-in controls / links so they aren't double-handled.
  document.addEventListener('click', ev => {{
    if (ev.target.closest('.controls, a, .navigate-left, .navigate-right, .navigate-up, .navigate-down')) return;
    Reveal.next();
  }});
</script>
</body>
</html>
"""


def build_web_export(pptx_path: Path, out_dir: Path, render_ps1: Path, reveal_src: Path) -> None:
    """Full pipeline: render PNGs, extract videos, vendor reveal.js, emit index.html."""
    out_dir.mkdir(parents=True, exist_ok=True)

    # Rebuild assets from scratch: assets are only ever written by name, so slides that
    # disappear between builds (hidden, deleted, renumbered) would linger as orphans.
    out_assets = out_dir / "assets"
    if out_assets.exists():
        shutil.rmtree(out_assets)
    out_assets.mkdir(parents=True)

    # Vendor reveal.js locally (no CDN -> fully self-contained site).
    shutil.copytree(reveal_src, out_assets / "reveal")

    hidden = hidden_slide_numbers(pptx_path)
    image_names = render_pngs(pptx_path, render_ps1, out_assets, hidden)
    slides = extract_slides(pptx_path, out_assets, image_names, hidden)

    title = pptx_path.stem.replace("_", " ")
    (out_dir / "index.html").write_text(render_index_html(slides, title), encoding="utf-8")
    logger.info(f"wrote {out_dir / 'index.html'} ({len(slides)} slides)")
