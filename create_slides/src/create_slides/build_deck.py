from pathlib import Path

from pptx import Presentation

CHAPTERS = [
    "Introduction",
    "Data Acquisition & Feature Engineering",
    "Generative Modeling for Molecular Discovery",
    "Predictive Modeling of Adsorption Capacity",
    "Polymer Filtering",
    "Untrained Model: Agglomerative Hierarchical Clustering",
    "Final Result: End-to-End Pipeline",
    "Conclusion and Future Work",
]

OUTPUT_PATH = Path(__file__).resolve().parents[2] / "output" / "thesis_slides.pptx"


def build_deck(output_path: Path = OUTPUT_PATH) -> Path:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Polymer Discovery for Wastewater Adsorption"
    title_slide.placeholders[1].text = "Master Thesis — Samuele Maggiori"

    section_layout = prs.slide_layouts[1]
    for chapter in CHAPTERS:
        slide = prs.slides.add_slide(section_layout)
        slide.shapes.title.text = chapter
        slide.placeholders[1].text_frame.text = "TODO: fill in content"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    saved_to = build_deck()
    print(f"Saved deck to {saved_to}")
